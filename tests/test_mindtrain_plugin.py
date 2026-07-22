import importlib.util
import json
import os
import stat
import subprocess
import sys
import tempfile
import unittest
import xml.etree.ElementTree as ET
from pathlib import Path
from unittest.mock import patch

import yaml


ROOT = Path(__file__).resolve().parents[1]
BRIDGE_PATH = ROOT / "plugins/mindtrain/scripts/mindtrain_mcp_bridge.py"
SPEC = importlib.util.spec_from_file_location("mindtrain_mcp_bridge", BRIDGE_PATH)
BRIDGE = importlib.util.module_from_spec(SPEC)
SPEC.loader.exec_module(BRIDGE)


class MindTrainPluginTest(unittest.TestCase):
    def test_release_versions_and_minimum_plugin_version_stay_aligned(self):
        namespace = {"m": "http://maven.apache.org/POM/4.0.0"}
        project_version = ET.parse(ROOT / "pom.xml").getroot().findtext("m:version", namespaces=namespace)
        manifest_version = json.loads(
            (ROOT / "plugins/mindtrain/.codex-plugin/plugin.json").read_text()
        )["version"]
        openapi_version = yaml.safe_load(
            (ROOT / "contracts/openapi/trainer-core.yaml").read_text()
        )["info"]["version"]
        compatibility_source = (
            ROOT / "apps/trainer-mcp/src/main/java/io/github/shigella520/mindtrain/mcp/McpCompatibility.java"
        ).read_text()

        expected = ".".join(str(part) for part in BRIDGE.normalized_version(manifest_version))
        self.assertEqual(expected, project_version.removesuffix("-SNAPSHOT"))
        self.assertEqual(expected, openapi_version)
        self.assertIn(f'MINIMUM_PLUGIN_VERSION = "{expected}"', compatibility_source)

    def test_marketplace_plugin_and_mcp_manifests_are_consistent(self):
        marketplace = json.loads((ROOT / ".agents/plugins/marketplace.json").read_text())
        plugin = json.loads((ROOT / "plugins/mindtrain/.codex-plugin/plugin.json").read_text())
        mcp = json.loads((ROOT / "plugins/mindtrain/.mcp.json").read_text())

        self.assertEqual("mindtrain", marketplace["name"])
        self.assertEqual("mindtrain", marketplace["plugins"][0]["name"])
        self.assertEqual("./plugins/mindtrain", marketplace["plugins"][0]["source"]["path"])
        self.assertEqual("ON_USE", marketplace["plugins"][0]["policy"]["authentication"])
        self.assertEqual("mindtrain", plugin["name"])
        self.assertEqual("./.mcp.json", plugin["mcpServers"])
        self.assertIn("mindtrain", mcp["mcpServers"])
        for field in ("composerIcon", "logo", "logoDark"):
            asset = ROOT / "plugins/mindtrain" / plugin["interface"][field]
            self.assertTrue(asset.is_file(), f"missing plugin asset: {asset}")

    def test_skill_and_starter_prompts_use_the_mindtrain_name(self):
        plugin = json.loads((ROOT / "plugins/mindtrain/.codex-plugin/plugin.json").read_text())
        prompts = plugin["interface"]["defaultPrompt"]

        self.assertTrue(all("$mindtrain" in prompt for prompt in prompts))
        self.assertTrue(all(prompt.isascii() for prompt in prompts))
        self.assertFalse((ROOT / "plugins/mindtrain/skills/knowledge-trainer").exists())
        self.assertFalse((ROOT / "skills/mindtrain").exists())

        skill_path = ROOT / "plugins/mindtrain/skills/mindtrain/SKILL.md"
        frontmatter = skill_path.read_text().split("---", 2)[1]
        self.assertEqual("mindtrain", yaml.safe_load(frontmatter)["name"])

    def test_bridge_exposes_configuration_and_training_tools(self):
        definitions = BRIDGE.tool_definitions()
        names = {definition["name"] for definition in definitions}
        self.assertIn("configure_mindtrain_instance", names)
        self.assertIn("get_mindtrain_configuration", names)
        self.assertEqual(BRIDGE.REMOTE_TOOL_NAMES, names - {
            "configure_mindtrain_instance", "get_mindtrain_configuration"
        } - BRIDGE.REFERENCE_TOOL_NAMES)
        self.assertIn("revise_saved_question", names)
        self.assertIn("reject_generated_question", names)
        self.assertIn("preview_knowledge_catalog_import", names)
        self.assertIn("list_knowledge_domains", names)
        self.assertIn("get_knowledge_catalog_tree", names)
        self.assertIn("search_knowledge_topics", names)
        self.assertIn("preview_training_domain", names)
        self.assertIn("confirm_training_domain", names)
        self.assertIn("sync_reference_library", names)
        session_tool = next(item for item in definitions if item["name"] == "create_training_session")
        self.assertNotIn("java-backend", json.dumps(session_tool))
        self.assertIn("multiple domains", json.dumps(session_tool))

    def test_skill_documents_dialogue_and_local_catalog_workflows(self):
        skill_root = ROOT / "plugins/mindtrain/skills/mindtrain"
        skill = (skill_root / "SKILL.md").read_text()
        catalog = (skill_root / "references/knowledge-catalog.md").read_text()
        local = (skill_root / "references/reference-library.md").read_text()

        self.assertIn("ai_dialogue", skill)
        self.assertIn("list_knowledge_domains", catalog)
        self.assertIn("multiple root topics", catalog)
        self.assertIn("confirm_training_domain", local)

    def test_url_requires_https_for_non_local_instances(self):
        self.assertEqual(
            "https://train.example.com/mcp",
            BRIDGE.normalize_url("https://train.example.com"),
        )
        self.assertEqual(
            "http://127.0.0.1:8787/mcp",
            BRIDGE.normalize_url("http://127.0.0.1:8787/mcp"),
        )
        with self.assertRaisesRegex(ValueError, "HTTPS is required"):
            BRIDGE.normalize_url("http://train.example.com/mcp")

    def test_configuration_is_saved_with_private_permissions_and_never_reported(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "mindtrain" / "plugin.json"
            with patch.dict(os.environ, {"MINDTRAIN_PLUGIN_CONFIG": str(path)}):
                compatibility = {
                    "status": "compatible", "pluginVersion": BRIDGE.PLUGIN_VERSION,
                    "serverVersion": "0.2.0-SNAPSHOT", "contractVersion": 1,
                    "minimumPluginVersion": "0.2.0", "versionMatch": True,
                }
                with patch.object(BRIDGE, "verify_config", return_value=compatibility):
                    result = BRIDGE.call_tool(
                        "configure_mindtrain_instance",
                        {"url": "https://train.example.com/mcp", "token": "secret-token"},
                    )
                    status = BRIDGE.configuration_status()

                self.assertFalse(result["isError"])
                self.assertEqual("compatible", result["structuredContent"]["compatibility"]["status"])
                self.assertEqual("secret-token", BRIDGE.load_config()["token"])
                self.assertTrue(status["configured"])
                self.assertEqual("compatible", status["compatibility"]["status"])
                self.assertNotIn("token", status)
                if os.name != "nt":
                    self.assertEqual(stat.S_IRUSR | stat.S_IWUSR, stat.S_IMODE(path.stat().st_mode))

    def test_plugin_and_server_compatibility_is_checked_both_ways(self):
        compatible = {
            "serverInfo": {"name": "mindtrain-trainer-mcp", "version": "0.2.1"},
            "_meta": {"mindtrainCompatibility": {
                "contractVersion": BRIDGE.CONTRACT_VERSION,
                "minimumPluginVersion": "0.2.0",
            }},
        }
        with patch.object(BRIDGE, "remote_request", return_value=compatible) as request:
            result = BRIDGE.verify_config({"url": "https://train.example.com/mcp", "token": "secret"})
        self.assertEqual("compatible_version_difference", result["status"])
        self.assertIn("warning", result)
        self.assertEqual(BRIDGE.PLUGIN_VERSION, request.call_args.args[2]["clientInfo"]["version"])

        incompatible = {
            **compatible,
            "_meta": {"mindtrainCompatibility": {
                "contractVersion": BRIDGE.CONTRACT_VERSION + 1,
                "minimumPluginVersion": "0.2.0",
            }},
        }
        with patch.object(BRIDGE, "remote_request", return_value=incompatible):
            with self.assertRaisesRegex(ValueError, "契约版本"):
                BRIDGE.verify_config({"url": "https://train.example.com/mcp", "token": "secret"})

    def test_remote_calls_send_plugin_and_contract_versions(self):
        class Response:
            def __enter__(self):
                return self

            def __exit__(self, *_):
                return False

            def read(self):
                return b'{"jsonrpc":"2.0","id":1,"result":{}}'

        captured = {}

        def open_request(request, timeout):
            captured["headers"] = dict(request.header_items())
            return Response()

        with patch.object(BRIDGE, "urlopen", side_effect=open_request):
            BRIDGE.remote_request({"url": "https://train.example.com/mcp", "token": "secret"}, "ping")

        self.assertEqual(BRIDGE.PLUGIN_VERSION, captured["headers"]["X-mindtrain-plugin-version"])
        self.assertEqual(str(BRIDGE.CONTRACT_VERSION), captured["headers"]["X-mindtrain-contract-version"])

    def test_local_reference_library_indexes_searches_and_stays_private(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            library = root / "notes"
            library.mkdir()
            (library / "guide.md").write_text(
                "# Raft consensus\n\nA leader replicates log entries to followers.", encoding="utf-8"
            )
            (library / "notes.txt").write_text("Quorum requires a majority.", encoding="utf-8")
            (library / ".secret.txt").write_text("do not index", encoding="utf-8")
            config = root / "config" / "plugin.json"
            cache = root / "cache"
            environment = {"MINDTRAIN_PLUGIN_CONFIG": str(config), "MINDTRAIN_PLUGIN_CACHE": str(cache)}
            with patch.dict(os.environ, environment), patch.object(
                BRIDGE, "reference_runtime_python", return_value=Path("/unused/python")
            ):
                configured = BRIDGE.call_tool("configure_reference_library", {
                    "libraryId": "distributed-notes", "path": str(library), "makeDefault": True
                })
                self.assertFalse(configured["isError"])
                synced = BRIDGE.call_tool("sync_reference_library", {"libraryId": "distributed-notes"})
                self.assertEqual(2, synced["structuredContent"]["documentCount"])
                matches = BRIDGE.call_tool("search_reference_library", {
                    "libraryId": "distributed-notes", "query": "Raft"
                })["structuredContent"]["matches"]
                self.assertEqual("guide.md", matches[0]["relativePath"])
                self.assertTrue(matches[0]["url"].startswith("mindtrain-local://distributed-notes/"))
                self.assertNotIn(str(library), json.dumps(matches))
                excerpt = BRIDGE.call_tool("read_reference_document", {
                    "libraryId": "distributed-notes", "relativePath": "guide.md", "maxCharacters": 30
                })["structuredContent"]
                self.assertLessEqual(len(excerpt["content"]), 30)
                rejected = BRIDGE.call_tool("read_reference_document", {
                    "libraryId": "distributed-notes", "relativePath": "../secret.txt"
                })
                self.assertTrue(rejected["isError"])

                stored = json.loads(config.read_text())
                self.assertEqual(str(library.resolve()), stored["referenceLibraries"][0]["path"])
                listed = BRIDGE.call_tool("list_reference_libraries", {})
                self.assertNotIn(str(library), json.dumps(listed))

    def test_reference_parser_lock_and_worker_are_packaged(self):
        scripts = ROOT / "plugins/mindtrain/scripts"
        lock = (scripts / "reference-requirements.lock").read_text()
        self.assertIn("pypdf==", lock)
        self.assertIn("python-docx==", lock)
        self.assertIn("python-pptx==", lock)
        self.assertTrue((scripts / "reference_worker.py").is_file())

    def test_reference_worker_extracts_pdf_docx_and_pptx(self):
        from docx import Document
        from pptx import Presentation
        from pypdf import PdfWriter

        worker = ROOT / "plugins/mindtrain/scripts/reference_worker.py"
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            docx_path = root / "guide.docx"
            document = Document()
            document.add_heading("DOCX knowledge", level=1)
            document.add_paragraph("Document body")
            document.save(docx_path)

            pptx_path = root / "slides.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "PPTX knowledge"
            presentation.save(pptx_path)

            pdf_path = root / "scan.pdf"
            writer = PdfWriter()
            writer.add_blank_page(width=200, height=200)
            with pdf_path.open("wb") as handle:
                writer.write(handle)

            extracted = {}
            for path in (docx_path, pptx_path, pdf_path):
                completed = subprocess.run(
                    [sys.executable, str(worker), str(path)], check=True,
                    stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True,
                )
                extracted[path.suffix] = json.loads(completed.stdout)

            self.assertIn("DOCX knowledge", extracted[".docx"]["text"])
            self.assertIn("PPTX knowledge", extracted[".pptx"]["text"])
            self.assertIn("OCR is not enabled", extracted[".pdf"]["warning"])


if __name__ == "__main__":
    unittest.main()
