#!/usr/bin/env python3
"""Isolated document text extractor used by the local MindTrain bridge.

The worker emits one JSON object on stdout and never sends document content over
the network. Its third-party dependencies live in a versioned private venv.
"""

import json
import sys
from pathlib import Path


def extract_pdf(path):
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception as error:
            raise ValueError("encrypted PDF is not supported") from error
    pages = []
    empty_pages = 0
    for number, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        empty_pages += int(not text.strip())
        pages.append("\n\n[Page {}]\n{}".format(number, text))
    warning = None
    if empty_pages:
        warning = "{} PDF page(s) contained no extractable text; OCR is not enabled".format(empty_pages)
    return path.stem, "".join(pages).strip(), warning


def extract_docx(path):
    from docx import Document

    document = Document(str(path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            paragraphs.append("\t".join(cell.text for cell in row.cells))
    title = next((paragraph.strip() for paragraph in paragraphs if paragraph.strip()), path.stem)
    return title, "\n".join(paragraphs), None


def extract_pptx(path):
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides = []
    first_line = None
    for number, slide in enumerate(presentation.slides, 1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
        if first_line is None and lines:
            first_line = lines[0].strip()
        slides.append("\n\n[Slide {}]\n{}".format(number, "\n".join(lines)))
    return first_line or path.stem, "".join(slides).strip(), None


def extract(path):
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    raise ValueError("unsupported document type: " + suffix)


def main():
    if len(sys.argv) != 2:
        raise ValueError("expected one document path")
    path = Path(sys.argv[1]).resolve(strict=True)
    title, text, warning = extract(path)
    print(json.dumps({"title": title, "text": text, "warning": warning}, ensure_ascii=False))


if __name__ == "__main__":
    try:
        main()
    except Exception as error:
        print(str(error), file=sys.stderr)
        raise SystemExit(1)
